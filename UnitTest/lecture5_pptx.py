from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def readPPT(filePath):
    return Presentation(filePath)

def writeNewPPT(filePath):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Yo, Python!'
    slide.placeholders[1].text = 'Yes it is really awsome'

    prs.save(filePath)

def writePPT1(filePath):
    prs = Presentation(filePath)

    first_slide = prs.slides[0]
    first_slide.shapes[0].text_frame.paragraphs[0].text = 'Hello!'

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    text_frame = slide.shapes[0].text_frame
    p = text_frame.paragraphs[0]
    p.text = 'This is a paragraph'

    dir, _ = os.path.split(filePath)
    prs.save(f'{dir}/new_ppt.pptx')

def writePPT2(filePath):
    prs = Presentation()

    layout = prs.slide_layouts[3]

    slide = prs.slides.add_slide(layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Two Content Slide'

    tf = body_shape.text_frame
    tf.text = 'This is line 1'

    p = tf.add_paragraph()
    p.text = 'Again a line 2'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'And this is line 3'
    p.level = 2

    prs.save(filePath)
    
def writePPT3(filePath):
    prs = Presentation()

    layout = prs.slide_layouts[6]
    
    slide = prs.slides.add_slide(layout)

    txtBox = slide.shapes.add_textbox(
        # x, y, 너비, 길이
        Inches(2), Inches(2), Inches(5), Inches(1)
    )

    tf = txtBox.text_frame
    tf.text = "Wow! I'm inside a textbox"
    p = tf.add_paragraph()
    p.text = 'adding a new text'
    p.font.bold = True
    p.font.italic = True
    p.font.size = Pt(30)

    prs.save(filePath)

def writeShapesPPT(filePath):
    prs = Presentation()

    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding Shapes'

    shape1 = shapes.add_shape(
        MSO_SHAPE.RECTANGULAR_CALLOUT,
        # x, y, 너비, 길이
        Inches(3.5), Inches(2), Inches(2), Inches(2)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0x1E, 0x90, 0xFF)
    shape1.fill.fore_color.brightness = 0.4
    shape1.text = 'See! There is home!'

    shape2 = shapes.add_shape(
        MSO_SHAPE.ACTION_BUTTON_HOME,
        Inches(3.5), Inches(5), Inches(2), Inches(2)
    )
    shape2.text = 'Home'
    
    prs.save(filePath)

def writeTable(filePath):
    prs = Presentation()

    layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(layout)
    shapes = slide.shapes
    shapes.title.text = 'Students Data'

    rows = 4; cols = 3
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(1.2)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)

    table.cell(0, 0).text = 'Sr. No.'
    table.cell(0, 1).text = 'Student Name'
    table.cell(0, 2).text = 'Student Id'

    students = {
        1: ['John', 115],
        2: ['Mary', 119],
        3: ['Alice', 101]
    }

    for i in range(len(students)):
        table.cell(i+1, 0).text = str(i+1)
        table.cell(i+1, 1).text = str(students[i+1][0])
        table.cell(i+1, 2).text = str(students[i+1][1])

    prs.save(filePath)